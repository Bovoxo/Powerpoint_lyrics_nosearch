import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import io
import re

def clean_lyrics(text):
    text = re.sub(r'\[.*?\]', '', text)
    text = re.sub(r'(You might also like|See upcoming.*|Get tickets.*|[A-Z][a-z]+ [A-Z][a-z]+)', '', text)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return lines

def create_pptx_from_lyrics(lines):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    for i in range(0, len(lines), 4):
        slide = prs.slides.add_slide(blank_slide_layout)
        box_width = Inches(10)
        box_height = Inches(6)
        box_left = (prs.slide_width - box_width) / 2
        box_top = Inches(0.75)
        text_box = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
        tf = text_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for line in lines[i:i+4]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(32)
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

st.title("🎤 Text-to-PowerPoint: Lyrics Slideshow Generator")
st.markdown("Vlož text písně z Genius.com. Každé 4 řádky budou na jednom slidu.")

user_input = st.text_area("Vlož text zde:", height=300)

if st.button("Vytvořit PowerPoint"):
    if user_input.strip():
        cleaned_lines = clean_lyrics(user_input)
        pptx_file = create_pptx_from_lyrics(cleaned_lines)
        st.success("Prezentace byla úspěšně vytvořena!")
        st.download_button("📥 Stáhnout prezentaci", data=pptx_file, file_name="lyrics_presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.warning("Prosím vlož text písně.")
